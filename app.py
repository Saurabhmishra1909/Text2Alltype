import tkinter as tk
from tkinter import ttk, messagebox, filedialog, colorchooser
from docx import Document
from docx.shared import Pt, RGBColor
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt as PptPt
from pptx.dml.color import RGBColor as PptRGBColor
import markdown
import os
import json
import csv
import xml.etree.ElementTree as ET
from PIL import Image, ImageDraw, ImageFont
import textwrap
import datetime

class TextConverter:
    def __init__(self, root):
        self.root = root
        self.root.title("Advanced Text Format Converter")
        
        self.font_family = tk.StringVar(value="Arial")
        self.font_size = tk.StringVar(value="12")
        self.text_color = "#000000"
        self.save_directory = os.getcwd()
        self.alignment = tk.StringVar(value="left")

        self.create_gui()

    def create_gui(self):
        self.main_frame = ttk.Frame(self.root, padding="10")
        self.main_frame.grid(row=0, column=0, sticky=(tk.W, tk.E, tk.N, tk.S))
        
        self.create_style_frame()
        
        self.text_label = ttk.Label(self.main_frame, text="Enter your text:")
        self.text_label.grid(row=2, column=0, sticky=tk.W, pady=5)
        
        self.text_input = tk.Text(self.main_frame, height=10, width=60)
        self.text_input.grid(row=3, column=0, columnspan=3, pady=5)
        
        self.filename_label = ttk.Label(self.main_frame, text="Output filename (without extension):")
        self.filename_label.grid(row=4, column=0, sticky=tk.W, pady=5)
        
        self.filename_entry = ttk.Entry(self.main_frame, width=40)
        self.filename_entry.grid(row=4, column=1, sticky=tk.W, pady=5)
        
        self.create_save_location_frame()
        self.create_conversion_buttons()

    def create_style_frame(self):
        style_frame = ttk.LabelFrame(self.main_frame, text="Style Options", padding="5")
        style_frame.grid(row=0, column=0, columnspan=3, sticky=(tk.W, tk.E), pady=5)
        
        ttk.Label(style_frame, text="Font:").grid(row=0, column=0, padx=5)
        fonts = ["Arial", "Times New Roman", "Courier", "Helvetica", "Verdana"]
        font_combo = ttk.Combobox(style_frame, textvariable=self.font_family, values=fonts)
        font_combo.grid(row=0, column=1, padx=5)
        
        ttk.Label(style_frame, text="Size:").grid(row=0, column=2, padx=5)
        sizes = [str(i) for i in range(8, 73, 2)]
        size_combo = ttk.Combobox(style_frame, textvariable=self.font_size, values=sizes, width=5)
        size_combo.grid(row=0, column=3, padx=5)
        
        ttk.Button(style_frame, text="Text Color", command=self.choose_color).grid(row=0, column=4, padx=5)
        
        ttk.Radiobutton(style_frame, text="Left", variable=self.alignment, value="left").grid(row=0, column=5, padx=5)
        ttk.Radiobutton(style_frame, text="Center", variable=self.alignment, value="center").grid(row=0, column=6, padx=5)
        ttk.Radiobutton(style_frame, text="Right", variable=self.alignment, value="right").grid(row=0, column=7, padx=5)

    def create_save_location_frame(self):
        save_frame = ttk.Frame(self.main_frame)
        save_frame.grid(row=5, column=0, columnspan=3, sticky=(tk.W, tk.E), pady=5)
        
        ttk.Label(save_frame, text="Save Location:").grid(row=0, column=0, sticky=tk.W)
        self.save_location_label = ttk.Label(save_frame, text=self.save_directory, wraplength=400)
        self.save_location_label.grid(row=0, column=1, sticky=tk.W, padx=5)
        ttk.Button(save_frame, text="Browse", command=self.browse_location).grid(row=0, column=2, padx=5)

    def create_conversion_buttons(self):
        button_frame = ttk.LabelFrame(self.main_frame, text="Convert To", padding="5")
        button_frame.grid(row=6, column=0, columnspan=3, sticky=(tk.W, tk.E), pady=5)
        
        ttk.Button(button_frame, text="PDF", command=self.convert_to_pdf).grid(row=0, column=0, padx=5, pady=5)
        ttk.Button(button_frame, text="DOCX", command=self.convert_to_docx).grid(row=0, column=1, padx=5, pady=5)
        ttk.Button(button_frame, text="PPT", command=self.convert_to_ppt).grid(row=0, column=2, padx=5, pady=5)
        ttk.Button(button_frame, text="HTML", command=self.convert_to_html).grid(row=0, column=3, padx=5, pady=5)
        ttk.Button(button_frame, text="Markdown", command=self.convert_to_markdown).grid(row=0, column=4, padx=5, pady=5)
        ttk.Button(button_frame, text="TXT", command=self.convert_to_txt).grid(row=1, column=0, padx=5, pady=5)
        ttk.Button(button_frame, text="CSV", command=self.convert_to_csv).grid(row=1, column=1, padx=5, pady=5)
        ttk.Button(button_frame, text="XML", command=self.convert_to_xml).grid(row=1, column=2, padx=5, pady=5)
        ttk.Button(button_frame, text="JSON", command=self.convert_to_json).grid(row=1, column=3, padx=5, pady=5)
        ttk.Button(button_frame, text="Image", command=self.convert_to_image).grid(row=1, column=4, padx=5, pady=5)    

    def convert_to_markdown(self):
        save_path = self.get_save_path("md")
        text = self.get_text()
        if not save_path or not text:
            return
        
        try:
            md_content = f"# Text Content\n\n{text}"  # Basic Markdown content
            with open(save_path, 'w', encoding='utf-8') as f:
                f.write(md_content)
            messagebox.showinfo("Success", f"Markdown file created: {save_path}")
        except Exception as e:
            messagebox.showerror("Error", f"Failed to create Markdown: {str(e)}")
               
    def create_style_frame(self):
        style_frame = ttk.LabelFrame(self.main_frame, text="Style Options", padding="5")
        style_frame.grid(row=0, column=0, columnspan=3, sticky=(tk.W, tk.E), pady=5)
        
        # Font family
        ttk.Label(style_frame, text="Font:").grid(row=0, column=0, padx=5)
        fonts = ["Arial", "Times New Roman", "Courier", "Helvetica", "Verdana"]
        font_combo = ttk.Combobox(style_frame, textvariable=self.font_family, values=fonts)
        font_combo.grid(row=0, column=1, padx=5)
        
        # Font size
        ttk.Label(style_frame, text="Size:").grid(row=0, column=2, padx=5)
        sizes = [str(i) for i in range(8, 73, 2)]
        size_combo = ttk.Combobox(style_frame, textvariable=self.font_size, values=sizes, width=5)
        size_combo.grid(row=0, column=3, padx=5)
        
        # Color picker
        ttk.Button(style_frame, text="Text Color", command=self.choose_color).grid(row=0, column=4, padx=5)
        
        # Text alignment
        self.alignment = tk.StringVar(value="left")
        ttk.Radiobutton(style_frame, text="Left", variable=self.alignment, value="left").grid(row=0, column=5, padx=5)
        ttk.Radiobutton(style_frame, text="Center", variable=self.alignment, value="center").grid(row=0, column=6, padx=5)
        ttk.Radiobutton(style_frame, text="Right", variable=self.alignment, value="right").grid(row=0, column=7, padx=5)

    def create_save_location_frame(self):
        save_frame = ttk.Frame(self.main_frame)
        save_frame.grid(row=5, column=0, columnspan=3, sticky=(tk.W, tk.E), pady=5)
        
        ttk.Label(save_frame, text="Save Location:").grid(row=0, column=0, sticky=tk.W)
        self.save_location_label = ttk.Label(save_frame, text=self.save_directory, wraplength=400)
        self.save_location_label.grid(row=0, column=1, sticky=tk.W, padx=5)
        ttk.Button(save_frame, text="Browse", command=self.browse_location).grid(row=0, column=2, padx=5)

    def create_conversion_buttons(self):
        button_frame = ttk.LabelFrame(self.main_frame, text="Convert To", padding="5")
        button_frame.grid(row=6, column=0, columnspan=3, sticky=(tk.W, tk.E), pady=5)
        
        # Document formats
        ttk.Button(button_frame, text="PDF", command=self.convert_to_pdf).grid(row=0, column=0, padx=5, pady=5)
        ttk.Button(button_frame, text="DOCX", command=self.convert_to_docx).grid(row=0, column=1, padx=5, pady=5)
        ttk.Button(button_frame, text="PPT", command=self.convert_to_ppt).grid(row=0, column=2, padx=5, pady=5)
        
        # Web formats
        ttk.Button(button_frame, text="HTML", command=self.convert_to_html).grid(row=0, column=3, padx=5, pady=5)
        ttk.Button(button_frame, text="Markdown", command=self.convert_to_markdown).grid(row=0, column=4, padx=5, pady=5)
        
        # Data formats
        ttk.Button(button_frame, text="TXT", command=self.convert_to_txt).grid(row=1, column=0, padx=5, pady=5)
        ttk.Button(button_frame, text="CSV", command=self.convert_to_csv).grid(row=1, column=1, padx=5, pady=5)
        ttk.Button(button_frame, text="XML", command=self.convert_to_xml).grid(row=1, column=2, padx=5, pady=5)
        ttk.Button(button_frame, text="JSON", command=self.convert_to_json).grid(row=1, column=3, padx=5, pady=5)
        
        # Image format
        ttk.Button(button_frame, text="Image", command=self.convert_to_image).grid(row=1, column=4, padx=5, pady=5)

    def choose_color(self):
        color = colorchooser.askcolor(title="Choose Text Color")
        if color[1]:
            self.text_color = color[1]

    def browse_location(self):
        directory = filedialog.askdirectory()
        if directory:
            self.save_directory = directory
            self.save_location_label.config(text=directory)

    def get_save_path(self, extension):
        filename = self.filename_entry.get().strip()
        if not filename:
            messagebox.showerror("Error", "Please enter a filename")
            return None
        return os.path.join(self.save_directory, f"{filename}.{extension}")

    def get_text(self):
        text = self.text_input.get("1.0", tk.END).strip()
        if not text:
            messagebox.showerror("Error", "Please enter some text")
            return None
        return text

    def convert_to_pdf(self):
        save_path = self.get_save_path("pdf")
        text = self.get_text()
        if not save_path or not text:
            return
        
        try:
            pdf = FPDF()
            pdf.add_page()
            pdf.set_font(self.font_family.get(), size=int(self.font_size.get()))
            
            # Convert hex color to RGB
            color = tuple(int(self.text_color.lstrip('#')[i:i+2], 16) for i in (0, 2, 4))
            pdf.set_text_color(*color)
            
            # Set alignment
            pdf.set_font(self.font_family.get(), size=int(self.font_size.get()))
            pdf.multi_cell(0, 10, txt=text, align=self.alignment.get())
            
            pdf.output(save_path)
            messagebox.showinfo("Success", f"PDF file created: {save_path}")
        except Exception as e:
            messagebox.showerror("Error", f"Failed to create PDF: {str(e)}")

    def convert_to_docx(self):
        save_path = self.get_save_path("docx")
        text = self.get_text()
        if not save_path or not text:
            return
        
        try:
            doc = Document()
            style = doc.styles['Normal']
            font = style.font
            font.name = self.font_family.get()
            font.size = Pt(int(self.font_size.get()))
            
            # Convert hex color to RGB
            color = tuple(int(self.text_color.lstrip('#')[i:i+2], 16) for i in (0, 2, 4))
            font.color.rgb = RGBColor(*color)
            
            paragraph = doc.add_paragraph(text)
            paragraph.alignment = {"left": 0, "center": 1, "right": 2}[self.alignment.get()]
            
            doc.save(save_path)
            messagebox.showinfo("Success", f"DOCX file created: {save_path}")
        except Exception as e:
            messagebox.showerror("Error", f"Failed to create DOCX: {str(e)}")

    def convert_to_ppt(self):
        save_path = self.get_save_path("pptx")
        text = self.get_text()
        if not save_path or not text:
            return
        
        try:
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            
            # Set title
            title = slide.shapes.title
            title.text = "Text Content"
            
            # Set content with formatting
            content = slide.shapes.placeholders[1]
            text_frame = content.text_frame
            
            paragraph = text_frame.paragraphs[0]
            paragraph.text = text
            paragraph.font.name = self.font_family.get()
            paragraph.font.size = PptPt(int(self.font_size.get()))
            
            # Convert hex color to RGB
            color = tuple(int(self.text_color.lstrip('#')[i:i+2], 16) for i in (0, 2, 4))
            paragraph.font.color.rgb = PptRGBColor(*color)
            
            paragraph.alignment = {"left": 1, "center": 2, "right": 3}[self.alignment.get()]
            
            prs.save(save_path)
            messagebox.showinfo("Success", f"PPTX file created: {save_path}")
        except Exception as e:
            messagebox.showerror("Error", f"Failed to create PPTX: {str(e)}")

    def convert_to_html(self):
        save_path = self.get_save_path("html")
        text = self.get_text()
        if not save_path or not text:
            return
        
        try:
            html_content = f"""
            <!DOCTYPE html>
            <html>
            <head>
                <style>
                    body {{
                        font-family: {self.font_family.get()};
                        font-size: {self.font_size.get()}px;
                        color: {self.text_color};
                        text-align: {self.alignment.get()};
                    }}
                </style>
            </head>
            <body>
                <p>{text}</p>
            </body>
            </html>
            """
            
            with open(save_path, 'w', encoding='utf-8') as f:
                f.write(html_content)
            
            messagebox.showinfo("Success", f"HTML file created: {save_path}")
        except Exception as e:
            messagebox.showerror("Error", f"Failed to create HTML: {str(e)}")

    def convert_to_json(self):
        save_path = self.get_save_path("json")
        text = self.get_text()
        if not save_path or not text:
            return

        try:
            json_content = {
                'content': text,
                'styling': {
                    'font_family': self.font_family.get(),
                    'font_size': self.font_size.get(),
                    'text_color': self.text_color,
                    'alignment': self.alignment.get(),
                },
                'metadata': {
                    'created_at': str(datetime.datetime.now()),
                    'word_count': len(text.split()),
                    'character_count': len(text)
                }
            }
            
            with open(save_path, 'w', encoding='utf-8') as f:
                json.dump(json_content, f, indent=4)
            messagebox.showinfo("Success", f"JSON file created: {save_path}")
        except Exception as e:
            messagebox.showerror("Error", f"Failed to create JSON: {str(e)}")
            
    def convert_to_txt(self):
        save_path = self.get_save_path("txt")
        text = self.get_text()
        if not save_path or not text:
            return
        
        try:
            with open(save_path, 'w', encoding='utf-8') as f:
                f.write(text)
            
            messagebox.showinfo("Success", f"TXT file created: {save_path}")
        except Exception as e:
            messagebox.showerror("Error", f"Failed to create TXT: {str(e)}")

    def convert_to_csv(self):
        save_path = self.get_save_path("csv")
        text = self.get_text()
        if not save_path or not text:
            return
        
        try:
            # Split text into rows
            rows = text.split('\n')
            with open(save_path, 'w', newline='', encoding='utf-8') as f:
                writer = csv.writer(f)
                for row in rows:
                    writer.writerow([row])
            
            messagebox.showinfo("Success", f"CSV file created: {save_path}")
        except Exception as e:
            messagebox.showerror("Error", f"Failed to create CSV: {str(e)}")

    def convert_to_xml(self):
        save_path = self.get_save_path("xml")
        text = self.get_text()
        if not save_path or not text:
            return
        
        try:
            root = ET.Element("document")
            content = ET.SubElement(root, "content")
            content.text = text
            
            tree = ET.ElementTree(root)
            tree.write(save_path, encoding='utf-8', xml_declaration=True)
            
            messagebox.showinfo("Success", f"XML file created: {save_path}")
        except Exception as e:
            messagebox.showerror("Error", "Failed to create XML: " + str(e))
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(rtf_content)
            messagebox.showinfo("Success", f"RTF file created: {output_path}")
        except Exception as e:
            messagebox.showerror("Error", f"Failed to create RTF: {str(e)}")

    def convert_to_json(self):
        save_path = self.get_save_path("json")
        text = self.get_text()
        if not save_path or not text:
            return

        try:
            json_content = {
                'content': text,
                'styling': {
                    'font_family': self.font_family.get(),
                    'font_size': self.font_size.get(),
                    'text_color': self.text_color,
                    'alignment': self.alignment.get(),
                },
                'metadata': {
                    'created_at': str(datetime.datetime.now()),
                    'word_count': len(text.split()),
                    'character_count': len(text)
                }
            }
            
            with open(save_path, 'w', encoding='utf-8') as f:
                json.dump(json_content, f, indent=4)
            messagebox.showinfo("Success", f"JSON file created: {save_path}")
        except Exception as e:
            messagebox.showerror("Error", f"Failed to create JSON: {str(e)}")
            
    def convert_to_image(self):
        save_path = self.get_save_path("png")
        text = self.get_text()
        if not save_path or not text:
            return

    try:
        width, height = 800, 400
        image = Image.new("RGB", (width, height), "white")  # Or "RGBA" for transparency
        draw = ImageDraw.Draw(image)

        try: # Font loading try-except block
            font = ImageFont.truetype(self.font_family.get()+".ttf", int(self.font_size.get()))
        except IOError:
            try:
                font = ImageFont.truetype("arial.ttf", int(self.font_size.get()))
            except IOError:
                font = ImageFont.load_default()
                messagebox.showwarning("Warning", "No fonts found. Using default font.") # Inform user

        color = tuple(int(self.text_color.lstrip('#')[i:i+2], 16) for i in (0, 2, 4))

        max_width = width - 100  # Adjust as needed
        wrapped_text = textwrap.fill(text, width=max_width) # Wrap the text

        text_width, text_height = draw.textsize(wrapped_text, font=font) # Size of wrapped text
        x = (width - text_width) / 2
        y = (height - text_height) / 2

        draw.text((x, y), wrapped_text, font=font, fill=color) # Draw wrapped text
        image.save(save_path)
        messagebox.showinfo("Success", f"Image file created: {save_path}")

    except IOError as e:  # Catch font errors more specifically
        messagebox.showerror("Error", f"Failed to load font: {str(e)}")
    except Exception as e:
        messagebox.showerror("Error", f"Failed to create Image: {str(e)}") 
        
def main():
    root = tk.Tk()
    app = TextConverter(root)
    root.mainloop()

if __name__ == "__main__":
    main()
